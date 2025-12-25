#!/usr/bin/env python3
"""
Generate a professional PowerPoint presentation for the UV Python Tool Guide.

Design specifications:
- Accent color: #007b85 (teal) for highlights, lines, and design details
- Background: pure white (#FFFFFF)
- Font: Roboto only
- Clean, modern, impressive design
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

from pptx.dml.color import RGBColor

# Design constants
ACCENT_COLOR = RGBColor(0x00, 0x7B, 0x85)  # #007b85
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0x66, 0x66, 0x66)
VERY_LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)

# Font settings
FONT_NAME = "Roboto"
FONT_NAME_MONO = "Roboto Mono"

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_font(run, size, bold=False, color=BLACK, font_name=FONT_NAME):
    """Set font properties for a run."""
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title_slide(prs):
    """Add the title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add accent line at top
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(0.15)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()

    # Add UV logo/icon placeholder (decorative shape)
    icon = slide.shapes.add_shape(
        MSO_SHAPE.HEXAGON,
        Inches(6.0), Inches(1.5),
        Inches(1.3), Inches(1.3)
    )
    icon.fill.solid()
    icon.fill.fore_color.rgb = ACCENT_COLOR
    icon.line.fill.background()

    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(11.7), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "UV Python Tool Guide"
    set_font(run, 54, bold=True, color=DARK_GRAY)

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(11.7), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "The Fast Python Package & Project Manager"
    set_font(run, 28, color=ACCENT_COLOR)

    # Date
    date_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "December 2025"
    set_font(run, 18, color=LIGHT_GRAY)

    # Bottom accent line
    line2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.5), Inches(6.8),
        Inches(2.3), Inches(0.05)
    )
    line2.fill.solid()
    line2.fill.fore_color.rgb = ACCENT_COLOR
    line2.line.fill.background()


def add_section_slide(prs, section_number, section_title):
    """Add a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Left accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(0.4), SLIDE_HEIGHT
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_COLOR
    bar.line.fill.background()

    # Section number
    num_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(2), Inches(1.5))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{section_number:02d}"
    set_font(run, 72, bold=True, color=ACCENT_COLOR)

    # Section title
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.0), Inches(11), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = section_title
    set_font(run, 44, bold=True, color=DARK_GRAY)


def add_content_slide(prs, title, bullet_points, code_snippet=None):
    """Add a content slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Top accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(0.5),
        Inches(1.5), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_font(run, 36, bold=True, color=DARK_GRAY)

    # Calculate content area based on whether there's code
    if code_snippet:
        content_width = Inches(6.0)
        content_left = Inches(0.8)
        code_left = Inches(7.2)
        code_width = Inches(5.3)
    else:
        content_width = Inches(11.7)
        content_left = Inches(0.8)

    # Bullet points
    content_box = slide.shapes.add_textbox(content_left, Inches(1.8), content_width, Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.level = 0
        p.space_before = Pt(12)
        p.space_after = Pt(6)

        # Add bullet
        run = p.add_run()
        run.text = "  " + point
        set_font(run, 20, color=DARK_GRAY)

    # Code snippet (if provided)
    if code_snippet:
        # Code background
        code_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            code_left, Inches(1.8),
            code_width, Inches(5.0)
        )
        code_bg.fill.solid()
        code_bg.fill.fore_color.rgb = VERY_LIGHT_GRAY
        code_bg.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

        # Code text
        code_box = slide.shapes.add_textbox(
            Inches(7.4), Inches(2.0),
            Inches(5.0), Inches(4.6)
        )
        tf = code_box.text_frame
        tf.word_wrap = True

        lines = code_snippet.strip().split('\n')
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            run = p.add_run()
            run.text = line
            set_font(run, 14, color=DARK_GRAY, font_name=FONT_NAME_MONO)


def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Top accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(0.5),
        Inches(1.5), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_font(run, 36, bold=True, color=DARK_GRAY)

    # Create table
    num_cols = len(headers)
    num_rows = len(rows) + 1  # +1 for header

    table_width = Inches(11.7)
    table_height = Inches(0.5) * num_rows

    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.8), Inches(1.8),
        table_width, table_height
    ).table

    # Style header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_COLOR

        tf = cell.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = header
        set_font(run, 16, bold=True, color=WHITE)

    # Style data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)

            # Alternating row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = VERY_LIGHT_GRAY

            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = cell_text
            set_font(run, 14, color=DARK_GRAY)


def add_two_column_slide(prs, title, left_title, left_points, right_title, right_points):
    """Add a slide with two columns."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Top accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(0.5),
        Inches(1.5), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()

    # Main Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_font(run, 36, bold=True, color=DARK_GRAY)

    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = left_title
    set_font(run, 22, bold=True, color=ACCENT_COLOR)

    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(left_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = "  " + point
        set_font(run, 18, color=DARK_GRAY)

    # Divider line
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.5), Inches(1.8),
        Inches(0.03), Inches(5.0)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    divider.line.fill.background()

    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(6.9), Inches(1.7), Inches(5.5), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = right_title
    set_font(run, 22, bold=True, color=ACCENT_COLOR)

    # Right column content
    right_box = slide.shapes.add_textbox(Inches(6.9), Inches(2.3), Inches(5.5), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(right_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = "  " + point
        set_font(run, 18, color=DARK_GRAY)


def add_closing_slide(prs):
    """Add the closing/thank you slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Full width accent bar at bottom
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(6.5),
        SLIDE_WIDTH, Inches(1.0)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_COLOR
    bar.line.fill.background()

    # Decorative hexagon
    icon = slide.shapes.add_shape(
        MSO_SHAPE.HEXAGON,
        Inches(6.0), Inches(1.5),
        Inches(1.3), Inches(1.3)
    )
    icon.fill.solid()
    icon.fill.fore_color.rgb = ACCENT_COLOR
    icon.line.fill.background()

    # Main message
    msg_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.7), Inches(1.0))
    tf = msg_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Start Using UV Today"
    set_font(run, 48, bold=True, color=DARK_GRAY)

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(11.7), Inches(0.6))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "curl -LsSf https://astral.sh/uv/install.sh | sh"
    set_font(run, 22, color=LIGHT_GRAY, font_name=FONT_NAME_MONO)

    # Documentation link
    link_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.5))
    tf = link_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "docs.astral.sh/uv"
    set_font(run, 18, color=ACCENT_COLOR)


def create_presentation():
    """Create the complete presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 1. Title Slide
    add_title_slide(prs)

    # 2. Table of Contents
    add_content_slide(prs, "Agenda", [
        "What is UV?",
        "Installation",
        "Python Version Management",
        "Virtual Environments",
        "Project Management",
        "Dependency Management",
        "Running Python Scripts",
        "Tools and uvx",
        "pip-Compatible Commands",
        "Building and Publishing",
        "Docker Integration",
        "Best Practices"
    ])

    # 3. What is UV? - Section
    add_section_slide(prs, 1, "What is UV?")

    # 4. What is UV? - Content
    add_content_slide(prs, "What is UV?", [
        "Extremely fast Python package and project manager",
        "Written in Rust for maximum performance",
        "10-100x faster than pip and other tools",
        "Single binary - no Python required to install",
        "Unified toolchain replacing multiple tools",
        "Cross-platform: Linux, macOS, Windows"
    ])

    # 5. UV Replaces Table
    add_table_slide(prs, "UV Replaces Multiple Tools",
        ["Traditional Tool", "UV Equivalent"],
        [
            ["pip", "uv pip install"],
            ["pip-tools", "uv pip compile"],
            ["pipx", "uvx / uv tool"],
            ["poetry / pipenv", "uv init, uv add, uv sync"],
            ["pyenv", "uv python"],
            ["virtualenv", "uv venv"]
        ]
    )

    # 6. Installation - Section
    add_section_slide(prs, 2, "Installation")

    # 7. Installation - Content
    add_content_slide(prs, "Installing UV", [
        "macOS / Linux:",
        "  curl -LsSf https://astral.sh/uv/install.sh | sh",
        "",
        "macOS with Homebrew:",
        "  brew install uv",
        "",
        "Windows PowerShell:",
        "  irm https://astral.sh/uv/install.ps1 | iex",
        "",
        "Shell Autocompletion:",
        "  eval \"$(uv --generate-shell-completion bash)\""
    ])

    # 8. Python Version Management - Section
    add_section_slide(prs, 3, "Python Version Management")

    # 9. Python Version Management - Content
    add_content_slide(prs, "Managing Python Versions", [
        "Install Python automatically",
        "Manage multiple versions",
        "No need for pyenv anymore",
        "Automatic version detection"
    ],
    code_snippet="""# Install latest Python
uv python install

# Install specific version
uv python install 3.12

# Install multiple versions
uv python install 3.11 3.12 3.13

# List available versions
uv python list

# Run with specific version
uv run --python 3.12 script.py""")

    # 10. Virtual Environments - Section
    add_section_slide(prs, 4, "Virtual Environments")

    # 11. Virtual Environments - Content
    add_content_slide(prs, "Virtual Environments", [
        "UV auto-discovers .venv directory",
        "No activation required with UV commands",
        "Automatic environment creation on first run",
        "Compatible with standard venv workflow"
    ],
    code_snippet="""# Create virtual environment
uv venv

# Create with specific Python
uv venv --python 3.12

# Create with custom path
uv venv my_env

# UV commands auto-use .venv
uv pip install requests
uv run script.py

# Manual activation (optional)
source .venv/bin/activate""")

    # 12. Project Management - Section
    add_section_slide(prs, 5, "Project Management")

    # 13. Project Management - Content
    add_content_slide(prs, "Creating Projects", [
        "Modern project workflow like Poetry/npm",
        "Automatic pyproject.toml generation",
        "Support for apps, libraries, and packages",
        "Built-in build system"
    ],
    code_snippet="""# Create basic application
uv init my-project
cd my-project

# Create a library
uv init --lib my-library

# Create packaged app
uv init --package my-app

# Project structure:
# my-project/
#   .python-version
#   .venv/
#   README.md
#   main.py
#   pyproject.toml""")

    # 14. Dependency Management - Section
    add_section_slide(prs, 6, "Dependency Management")

    # 15. Dependency Management - Content
    add_content_slide(prs, "Managing Dependencies", [
        "Simple add/remove workflow",
        "Automatic lockfile management",
        "Support for dev and optional deps",
        "Git, local, and URL sources"
    ],
    code_snippet="""# Add dependencies
uv add requests flask

# Add with version constraint
uv add 'django>=4.0,<5.0'

# Add dev dependencies
uv add --dev pytest black mypy

# Remove dependencies
uv remove requests

# Lock and sync
uv lock
uv sync""")

    # 16. Dependency Sources
    add_content_slide(prs, "Dependency Sources", [
        "PyPI packages (default)",
        "Git repositories with branches/tags",
        "Local packages (editable installs)",
        "Direct URLs to archives",
        "Platform-specific dependencies"
    ],
    code_snippet="""# From Git
uv add git+https://github.com/org/repo
uv add git+https://github.com/org/repo --tag v1.0

# From local path
uv add ./local-package
uv add --editable ./my-lib

# Platform-specific
uv add "jax; sys_platform == 'linux'"

# Migrate from requirements.txt
uv add -r requirements.txt""")

    # 17. Running Scripts - Section
    add_section_slide(prs, 7, "Running Python Scripts")

    # 18. Running Scripts - Content
    add_content_slide(prs, "Running Python Scripts", [
        "Run scripts in project context",
        "Automatic dependency installation",
        "Inline script metadata support",
        "Ad-hoc dependencies for one-off runs"
    ],
    code_snippet="""# Run script in project
uv run script.py

# Run with arguments
uv run script.py --arg value

# Run Python module
uv run -m pytest

# Run with specific Python
uv run --python 3.12 script.py

# Ad-hoc dependencies
uv run --with requests fetch.py""")

    # 19. Inline Metadata
    add_content_slide(prs, "Inline Script Metadata", [
        "Embed dependencies directly in scripts",
        "Self-contained executable scripts",
        "Automatic dependency resolution",
        "Specify Python version requirements"
    ],
    code_snippet="""# /// script
# requires-python = ">=3.12"
# dependencies = [
#     "requests>=2.28",
#     "rich",
# ]
# ///

import requests
from rich import print

resp = requests.get("https://api.github.com")
print(resp.json())

# Then just run:
# uv run script.py""")

    # 20. Tools and uvx - Section
    add_section_slide(prs, 8, "Tools and uvx")

    # 21. Tools and uvx - Content
    add_two_column_slide(prs, "Tools and uvx",
        "Run Tools Temporarily (uvx)",
        [
            "uvx ruff check .",
            "uvx black --check .",
            "uvx mypy src/",
            "uvx ruff@0.6.0 --version",
            "uvx --from httpie http GET url"
        ],
        "Install Tools Permanently",
        [
            "uv tool install ruff",
            "uv tool install black",
            "uv tool list",
            "uv tool upgrade ruff",
            "uv tool uninstall ruff"
        ]
    )

    # 22. pip-Compatible Commands - Section
    add_section_slide(prs, 9, "pip-Compatible Commands")

    # 23. pip-Compatible Commands - Content
    add_content_slide(prs, "pip-Compatible Commands", [
        "Familiar pip interface",
        "Drop-in replacement for most pip commands",
        "Same flags and options",
        "10-100x faster execution"
    ],
    code_snippet="""# Install packages
uv pip install requests flask

# Install from requirements
uv pip install -r requirements.txt

# Editable install
uv pip install -e .

# List packages
uv pip list

# Compile requirements
uv pip compile requirements.in -o requirements.txt

# Freeze installed packages
uv pip freeze > requirements.txt""")

    # 24. Building and Publishing - Section
    add_section_slide(prs, 10, "Building and Publishing")

    # 25. Building and Publishing - Content
    add_content_slide(prs, "Building and Publishing Packages", [
        "Build source distributions and wheels",
        "Publish to PyPI or private registries",
        "Token-based authentication",
        "Multi-package workspace support"
    ],
    code_snippet="""# Build package
uv build

# Build only wheel
uv build --wheel

# Build only sdist
uv build --sdist

# Publish to PyPI
uv publish

# Publish with token
uv publish --token $PYPI_TOKEN

# Publish to custom registry
uv publish --publish-url https://my-registry/""")

    # 26. Docker Integration - Section
    add_section_slide(prs, 11, "Docker Integration")

    # 27. Docker Integration - Content
    add_content_slide(prs, "Docker Integration", [
        "Official UV Docker image available",
        "Optimized multi-stage builds",
        "Frozen lockfile for reproducibility",
        "Minimal image sizes"
    ],
    code_snippet="""FROM python:3.12-slim

# Install uv
COPY --from=ghcr.io/astral-sh/uv:latest \\
     /uv /bin/uv

WORKDIR /app

# Copy and install dependencies
COPY pyproject.toml uv.lock ./
RUN uv sync --frozen --no-dev

# Copy application
COPY . .

CMD ["uv", "run", "python", "main.py"]""")

    # 28. Best Practices - Section
    add_section_slide(prs, 12, "Best Practices")

    # 29. Best Practices - Content
    add_two_column_slide(prs, "Best Practices",
        "Do",
        [
            "Always use virtual environments",
            "Pin Python version (.python-version)",
            "Commit uv.lock to version control",
            "Use dependency groups",
            "Use --frozen in CI/CD",
            "Use uvx for one-off tools"
        ],
        "Avoid",
        [
            "Installing to system Python",
            "Ignoring the lockfile",
            "Mixing pip and uv in same project",
            "Manual dependency management",
            "Forgetting to sync after changes"
        ]
    )

    # 30. Quick Reference - Command Table
    add_table_slide(prs, "Quick Command Reference",
        ["Command", "Description"],
        [
            ["uv init", "Initialize new project"],
            ["uv add <pkg>", "Add dependency"],
            ["uv add --dev <pkg>", "Add dev dependency"],
            ["uv remove <pkg>", "Remove dependency"],
            ["uv lock", "Create/update lockfile"],
            ["uv sync", "Sync environment to lockfile"],
            ["uv run <cmd>", "Run command in project context"],
            ["uvx <tool>", "Run tool temporarily"],
            ["uv python install", "Install Python version"],
            ["uv venv", "Create virtual environment"],
            ["uv build", "Build package"],
            ["uv publish", "Publish to PyPI"]
        ]
    )

    # 31. Summary
    add_content_slide(prs, "Summary: UV Workflow", [
        "1. Install UV",
        "     curl -LsSf https://astral.sh/uv/install.sh | sh",
        "",
        "2. Create Project",
        "     uv init my-project && cd my-project",
        "",
        "3. Add Dependencies",
        "     uv add requests flask sqlalchemy",
        "",
        "4. Run Code",
        "     uv run python main.py",
        "",
        "5. Sync Environment",
        "     uv sync"
    ])

    # 32. Closing Slide
    add_closing_slide(prs)

    return prs


def main():
    """Main entry point."""
    print("Generating UV Python Tool Guide presentation...")

    prs = create_presentation()

    # Save the presentation
    output_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    output_path = os.path.join(output_dir, "UV-Python-Tool-Guide-Presentation.pptx")

    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
